#!/usr/bin/env python3
"""
Extrae contenido de archivos PPTX a texto plano con formato Markdown.
Preserva textos de diapositivas, notas del orador y referencias a imágenes.
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_pptx(file_path: str) -> dict:
    """
    Extrae el contenido de un archivo PPTX.

    Retorna:
        {
            'title': str,
            'slides': list[dict],  # Cada slide: {'number': int, 'title': str, 'text': str, 'notes': str}
            'error': str or None
        }
    """
    if Presentation is None:
        return {
            'title': Path(file_path).stem,
            'slides': [],
            'error': 'python-pptx no está instalado. Instala con: pip install python-pptx'
        }

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {
            'title': Path(file_path).stem,
            'slides': [],
            'error': f'Error al abrir PPTX: {str(e)}'
        }

    title = Path(file_path).stem
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'number': slide_num,
            'title': '',
            'text': [],
            'notes': ''
        }

        # Extrae texto de shapes (títulos, cajas de texto)
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                # Considera el primer shape como título si no tenemos uno
                if not slide_data['title'] and hasattr(shape, 'name') and 'Title' in shape.name:
                    slide_data['title'] = text
                else:
                    slide_data['text'].append(text)

        # Extrae notas del orador
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data['notes'] = notes_text

        slides.append(slide_data)

    return {
        'title': title,
        'slides': slides,
        'error': None
    }


def format_for_markdown(extracted: dict) -> str:
    """
    Formatea el contenido extraído como Markdown legible.
    """
    if extracted['error']:
        return f"# Error\n\n{extracted['error']}"

    output = [f"# {extracted['title']}\n"]

    for slide in extracted['slides']:
        output.append(f"## Diapositiva {slide['number']}")
        if slide['title']:
            output.append(f"### {slide['title']}")
        for text_block in slide['text']:
            output.append(f"\n{text_block}")
        if slide['notes']:
            output.append(f"\n*Notas del orador:*\n{slide['notes']}")
        output.append('')

    return '\n'.join(output)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({
            'error': 'Uso: extract_pptx.py <ruta-pptx>',
            'slides': [],
            'title': ''
        }))
        sys.exit(1)

    file_path = sys.argv[1]
    result = extract_pptx(file_path)

    if '--json' in sys.argv:
        print(json.dumps(result))
    else:
        print(format_for_markdown(result))
